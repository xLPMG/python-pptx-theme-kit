"""Primitive drawing helpers for python-pptx slides."""

import os

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


EMU_PER_POINT = 12700


def make_primitives(palette):
    """Build and return a dictionary of themed slide helper primitives.

    Args:
        palette: Dictionary returned by ``get_palette`` with color keys used by
            the drawing helpers.

    Returns:
        Dict[str, Callable]: Primitive helper functions for composing slides.
    """
    dark_bg = palette["DARK_BG"]
    code_bg = palette["CODE_BG"]
    accent = palette["ACCENT"]
    code_fg = palette["ACCENT2"]
    accent3 = palette["ACCENT3"]
    white = palette["WHITE"]
    light_grey = palette["LIGHT_GREY"]
    subtitle_c = palette["SUBTITLE_C"]
    row_a = palette["ROW_A"]
    row_b = palette["ROW_B"]

    slide_w = Inches(13.33)

    def set_bg(slide, color):
        """Set a solid background color on a slide.

        Args:
            slide: Target slide object.
            color: RGBColor value for the background.

        Returns:
            None.
        """
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(
        slide,
        left,
        top,
        width,
        height,
        fill_color,
        line_color=None,
        line_width=None,
    ):
        """Add a filled rectangle with an optional border.

        Args:
            slide: Target slide object.
            left: Left position in EMU/Inches.
            top: Top position in EMU/Inches.
            width: Rectangle width in EMU/Inches.
            height: Rectangle height in EMU/Inches.
            fill_color: RGBColor value for rectangle fill.
            line_color: Optional RGBColor border color.
            line_width: Optional border width (Pt/EMU).

        Returns:
            The created rectangle shape.
        """
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            if line_width:
                shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape

    def add_text(
        slide,
        text,
        left,
        top,
        width,
        height,
        size=14,
        bold=False,
        italic=False,
        color=white,
        align=PP_ALIGN.LEFT,
        wrap=True,
    ):
        """Add a text box with one formatted run.

        Args:
            slide: Target slide object.
            text: Text content to render.
            left: Left position in EMU/Inches.
            top: Top position in EMU/Inches.
            width: Text box width in EMU/Inches.
            height: Text box height in EMU/Inches.
            size: Font size in points.
            bold: Whether text is bold.
            italic: Whether text is italic.
            color: RGBColor value for font color.
            align: Paragraph alignment from PP_ALIGN.
            wrap: Whether to enable word wrapping.

        Returns:
            The created text box shape.
        """
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        return tb

    def add_image(
        slide,
        path,
        left,
        top,
        width=None,
        height=None,
        fit="contain",
        align="center",
        border_color=None,
        border_width=None,
    ):
        """Add an image with configurable sizing behavior.

        Supported fit modes:
            - ``native``: Use image's native size unless width/height is given.
            - ``contain``: Preserve aspect ratio and fit fully within frame.
            - ``cover``: Fill frame while preserving aspect ratio via cropping.
            - ``stretch``: Fill frame exactly (may distort aspect ratio).

        Args:
            slide: Target slide object.
            path: Image path on disk.
            left: Left position in EMU/Inches.
            top: Top position in EMU/Inches.
            width: Frame width in EMU/Inches.
            height: Frame height in EMU/Inches.
            fit: Image fit mode.
            align: Alignment hint for ``contain``/``cover``.
            border_color: Optional RGBColor border color.
            border_width: Optional border width (Pt/EMU).

        Returns:
            The created image shape.
        """
        if not os.path.exists(path):
            raise FileNotFoundError(f"Image file not found: {path}")

        fit_mode = (fit or "contain").lower()
        if fit_mode not in {"native", "contain", "cover", "stretch"}:
            raise ValueError(f"Unsupported fit mode: {fit}")

        if width is None and height is None:
            pic = slide.shapes.add_picture(path, left, top)
        elif width is None:
            pic = slide.shapes.add_picture(path, left, top, height=height)
        elif height is None:
            pic = slide.shapes.add_picture(path, left, top, width=width)
        elif fit_mode == "stretch":
            pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
        elif fit_mode == "cover":
            pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
            native = slide.shapes.add_picture(path, 0, 0)
            image_ratio = native.width / native.height
            native._element.getparent().remove(native._element)

            frame_ratio = width / height
            h_align = "center"
            v_align = "center"
            align_norm = (align or "center").lower().replace("_", "-")
            if "left" in align_norm:
                h_align = "left"
            elif "right" in align_norm:
                h_align = "right"
            if "top" in align_norm:
                v_align = "top"
            elif "bottom" in align_norm:
                v_align = "bottom"

            if image_ratio > frame_ratio:
                visible = frame_ratio / image_ratio
                if h_align == "left":
                    pic.crop_left = 0
                    pic.crop_right = 1 - visible
                elif h_align == "right":
                    pic.crop_left = 1 - visible
                    pic.crop_right = 0
                else:
                    crop = (1 - visible) / 2
                    pic.crop_left = crop
                    pic.crop_right = crop
            elif image_ratio < frame_ratio:
                visible = image_ratio / frame_ratio
                if v_align == "top":
                    pic.crop_top = 0
                    pic.crop_bottom = 1 - visible
                elif v_align == "bottom":
                    pic.crop_top = 1 - visible
                    pic.crop_bottom = 0
                else:
                    crop = (1 - visible) / 2
                    pic.crop_top = crop
                    pic.crop_bottom = crop
        elif fit_mode == "contain" and width is not None and height is not None:
            pic = slide.shapes.add_picture(path, 0, 0)
            native_w = pic.width
            native_h = pic.height
            scale = min(width / native_w, height / native_h)
            draw_w = int(native_w * scale)
            draw_h = int(native_h * scale)

            h_align = "center"
            v_align = "center"
            align_norm = (align or "center").lower().replace("_", "-")
            if "left" in align_norm:
                h_align = "left"
            elif "right" in align_norm:
                h_align = "right"
            if "top" in align_norm:
                v_align = "top"
            elif "bottom" in align_norm:
                v_align = "bottom"

            if h_align == "left":
                offset_x = 0
            elif h_align == "right":
                offset_x = width - draw_w
            else:
                offset_x = (width - draw_w) / 2

            if v_align == "top":
                offset_y = 0
            elif v_align == "bottom":
                offset_y = height - draw_h
            else:
                offset_y = (height - draw_h) / 2

            pic.left = int(left + offset_x)
            pic.top = int(top + offset_y)
            pic.width = draw_w
            pic.height = draw_h
        else:
            pic = slide.shapes.add_picture(path, left, top, width=width, height=height)

        if border_color:
            pic.line.color.rgb = border_color
            if border_width:
                pic.line.width = border_width

        return pic

    def add_code(
        slide,
        code,
        left,
        top,
        width,
        height,
        size=10,
        min_size=8,
        wrap=False,
        truncate=True,
    ):
        """Add a code panel with auto-fit font sizing and optional truncation.

        The function draws a themed code container, then inserts code lines into
        an inner text box. Font size is reduced down to ``min_size`` when needed
        to fit vertically. If lines still do not fit and ``truncate`` is True,
        remaining lines are dropped and the final rendered line gets ``...``.

        Args:
            slide: Target slide object.
            code: Multi-line code string to render.
            left: Left position in EMU/Inches.
            top: Top position in EMU/Inches.
            width: Code panel width in EMU/Inches.
            height: Code panel height in EMU/Inches.
            size: Preferred font size in points.
            min_size: Minimum font size allowed during fit.
            wrap: Whether to enable word wrapping in code text.
            truncate: Whether to truncate lines that exceed vertical space.

        Returns:
            None.
        """
        add_rect(slide, left, top, width, height, code_bg, accent, Pt(1))
        tb = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.1),
            width - Inches(0.3), height - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = wrap

        lines = code.split("\n")
        if not lines:
            lines = [""]

        available_h = int(height - Inches(0.2))
        max_size = int(size)
        min_size = int(min_size)
        line_spacing = 1.35
        fit_size = max_size

        while fit_size > min_size:
            required_h = int(len(lines) * fit_size * line_spacing * EMU_PER_POINT)
            if required_h <= available_h:
                break
            fit_size -= 1

        max_lines = int(available_h / (fit_size * line_spacing * EMU_PER_POINT))
        if max_lines < 1:
            max_lines = 1
        if truncate and len(lines) > max_lines:
            lines = lines[:max_lines]
            lines[-1] = f"{lines[-1]} ..."

        first = True
        for line in lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = line
            r.font.size = Pt(fit_size)
            r.font.name = "Courier New"
            r.font.color.rgb = code_fg

    def blank_slide(prs):
        """Create a blank layout slide (layout index 6).

        Args:
            prs: Presentation object.

        Returns:
            The newly created slide.
        """
        return prs.slides.add_slide(prs.slide_layouts[6])

    def title_bar(
        slide,
        title,
        subtitle="",
        bar_height=Inches(1.05),
    ):
        """Draw a full-width top bar with a title and optional subtitle.

        Args:
            slide: Target slide object.
            title: Main title text.
            subtitle: Optional subtitle text.
            bar_height: Height of the title bar.

        Returns:
            None.
        """
        add_rect(slide, 0, 0, slide_w, bar_height, accent)
        add_text(slide, title, Inches(0.35), Inches(0.1),
                 Inches(12.6), Inches(0.58), size=26, bold=True,
                 color=dark_bg)
        if subtitle:
            add_text(slide, subtitle, Inches(0.35), Inches(0.68),
                     Inches(12.6), Inches(0.32), size=13,
                     italic=True, color=dark_bg)

    def section_label(
        slide,
        text,
        left,
        top,
        width,
        height=Inches(0.3),
        size=13,
        color=accent,
        align=PP_ALIGN.LEFT,
    ):
        """Add a bold section label scoped to the provided rectangular region.

        Args:
            slide: Target slide object.
            text: Label text.
            left: Left position in EMU/Inches.
            top: Top position in EMU/Inches.
            width: Label width in EMU/Inches.
            height: Label height in EMU/Inches.
            size: Font size in points.
            color: RGBColor value for font color.
            align: Paragraph alignment from PP_ALIGN.

        Returns:
            None.
        """
        add_text(
            slide,
            text,
            left,
            top,
            width,
            height,
            size=size,
            bold=True,
            color=color,
            align=align,
        )

    def bullet_block(
        slide,
        items,
        left,
        top,
        width,
        height,
        size=13,
        color=light_grey,
        bullet="•",
    ):
        """Render a vertical bullet list inside a text box.

        Args:
            slide: Target slide object.
            items: Iterable of bullet item strings.
            left: Left position in EMU/Inches.
            top: Top position in EMU/Inches.
            width: Text box width in EMU/Inches.
            height: Text box height in EMU/Inches.
            size: Font size in points.
            color: RGBColor value for font color.
            bullet: Bullet marker prefix.

        Returns:
            None.
        """
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            r = p.add_run()
            r.text = f"{bullet}  {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = color

    def info_row(
        slide,
        label,
        value,
        y,
        row_h=Inches(0.52),
        lw=Inches(2.8),
        label_color=accent3,
        value_color=light_grey,
    ):
        """Draw a striped key/value row used for compact metrics summaries.

        Args:
            slide: Target slide object.
            label: Left-column label text.
            value: Right-column value text.
            y: Vertical position of the row.
            row_h: Row height.
            lw: Label-column width.
            label_color: RGBColor for label text.
            value_color: RGBColor for value text.

        Returns:
            None.
        """
        bg = row_a if int(y / Inches(0.52)) % 2 == 0 else row_b
        add_rect(slide, Inches(0.35), y, Inches(12.6), row_h, bg)
        add_text(slide, label, Inches(0.5), y + Inches(0.08),
                 lw - Inches(0.2), row_h - Inches(0.1),
                 size=12, bold=True, color=label_color)
        add_text(slide, value, Inches(0.5) + lw, y + Inches(0.08),
                 Inches(12.6) - lw - Inches(0.3), row_h - Inches(0.1),
                 size=12, color=value_color)

    def footer(slide, text):
        """Add centered footer text at the bottom of the slide.

        Args:
            slide: Target slide object.
            text: Footer text content.

        Returns:
            None.
        """
        add_text(slide, text, Inches(0.35), Inches(7.1),
                 Inches(12.6), Inches(0.3), size=10,
                 italic=True, color=subtitle_c, align=PP_ALIGN.CENTER)

    return {
        "set_bg": set_bg,
        "add_rect": add_rect,
        "add_text": add_text,
        "add_image": add_image,
        "add_code": add_code,
        "title_bar": title_bar,
        "section_label": section_label,
        "bullet_block": bullet_block,
        "info_row": info_row,
        "blank_slide": blank_slide,
        "footer": footer,
    }
